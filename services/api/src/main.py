"""FastAPI gateway for proposal-diagram-copilot"""
import os
import uuid
import asyncio
import logging
from datetime import datetime
from typing import Optional
from pathlib import Path

from fastapi import FastAPI, HTTPException, UploadFile, File, Form
from fastapi.middleware.cors import CORSMiddleware
from fastapi.responses import FileResponse
from pydantic import BaseModel

from shared_types import (
    JobStatus, OutputFormat, Job, Asset,
    UploadSessionRequest, UploadSessionResponse,
    CreateJobRequest, CreateJobResponse,
    JobStatusResponse, JobResult, QualityScore
)

# Setup logging
logging.basicConfig(
    level=logging.INFO,
    format="%(asctime)s [%(levelname)s] %(name)s: %(message)s"
)
logger = logging.getLogger("api")

app = FastAPI(
    title="Proposal Diagram Copilot API",
    version="0.1.0",
    docs_url="/docs",
    redoc_url="/redoc"
)

# CORS for Next.js
app.add_middleware(
    CORSMiddleware,
    allow_origins=["http://localhost:3000", "http://127.0.0.1:3000"],
    allow_credentials=True,
    allow_methods=["*"],
    allow_headers=["*"],
)

# ============ Storage (In-memory fallback) ============
class MemoryStore:
    """内存存储，无 Redis 时的 fallback"""
    def __init__(self):
        self.jobs: dict[str, Job] = {}
        self.assets: dict[str, Asset] = {}
        self.queue: list[str] = []  # job_ids
        
    def add_job(self, job: Job):
        self.jobs[job.id] = job
        self.queue.append(job.id)
        logger.info(f"[Store] Job {job.id} added to queue")
        
    def get_job(self, job_id: str) -> Optional[Job]:
        return self.jobs.get(job_id)
    
    def update_job(self, job: Job):
        job.updated_at = datetime.now()
        self.jobs[job.id] = job
        logger.info(f"[Store] Job {job.id} updated: status={job.status}, progress={job.progress}%")
        
    def add_asset(self, asset: Asset):
        self.assets[asset.id] = asset
        logger.info(f"[Store] Asset {asset.id} created: {asset.filename}")
        
    def get_asset(self, asset_id: str) -> Optional[Asset]:
        return self.assets.get(asset_id)
    
    def pop_queue(self) -> Optional[str]:
        if self.queue:
            return self.queue.pop(0)
        return None

store = MemoryStore()

# ============ Redis (Optional) ============
try:
    import redis as redis_lib
    REDIS_URL = os.getenv("REDIS_URL", "redis://localhost:6379/0")
    redis_client = redis_lib.from_url(REDIS_URL)
    redis_client.ping()
    USE_REDIS = True
    logger.info(f"[Redis] Connected to {REDIS_URL}")
except Exception as e:
    USE_REDIS = False
    logger.warning(f"[Redis] Not available, using memory fallback: {e}")

# ============ File Storage ============
UPLOAD_DIR = Path(__file__).parent / "uploads"
OUTPUT_DIR = Path(__file__).parent / "outputs"
UPLOAD_DIR.mkdir(exist_ok=True)
OUTPUT_DIR.mkdir(exist_ok=True)

# ============ Health Check ============
@app.get("/v1/health")
async def health():
    return {
        "status": "ok",
        "redis": USE_REDIS,
        "timestamp": datetime.now().isoformat()
    }

# ============ Asset Upload ============
@app.post("/v1/assets/upload-session", response_model=UploadSessionResponse)
async def create_upload_session(req: UploadSessionRequest):
    """创建上传会话（MVP 简化：直接返回 asset_id，实际上传用 /v1/assets/upload）"""
    asset_id = f"a_{uuid.uuid4().hex[:12]}"
    asset = Asset(
        id=asset_id,
        project_id=req.project_id,
        filename=req.filename,
        content_type=req.content_type,
        size=req.size,
        path="",  # 填充实际上传后
        created_at=datetime.now()
    )
    store.add_asset(asset)
    
    # MVP: 返回一个 fake upload URL（实际用 /v1/assets/upload）
    upload_url = f"http://localhost:8000/v1/assets/upload/{asset_id}"
    
    return UploadSessionResponse(
        asset_id=asset_id,
        upload_url=upload_url,
        expire_at=datetime.now()
    )

@app.post("/v1/assets/upload/{asset_id}")
async def upload_file(asset_id: str, file: UploadFile = File(...)):
    """实际上传文件"""
    asset = store.get_asset(asset_id)
    if not asset:
        raise HTTPException(404, "Asset not found")
    
    # 保存文件
    file_path = UPLOAD_DIR / f"{asset_id}_{file.filename}"
    content = await file.read()
    with open(file_path, "wb") as f:
        f.write(content)
    
    # 更新 asset
    asset.path = str(file_path)
    store.assets[asset_id] = asset
    
    logger.info(f"[Upload] File saved: {file_path} ({len(content)} bytes)")
    return {"asset_id": asset_id, "size": len(content), "path": str(file_path)}

# ============ Jobs ============
@app.post("/v1/jobs", response_model=CreateJobResponse)
async def create_job(req: CreateJobRequest):
    """创建转换任务"""
    asset = store.get_asset(req.asset_id)
    if not asset:
        raise HTTPException(404, "Asset not found")
    
    job_id = f"j_{uuid.uuid4().hex[:12]}"
    job = Job(
        id=job_id,
        project_id=req.project_id,
        asset_id=req.asset_id,
        status=JobStatus.QUEUED,
        stage="queued",
        progress=0,
        output_format=req.output_format,
        template_id=req.template_id,
        options=req.options,
        created_at=datetime.now(),
        updated_at=datetime.now()
    )
    store.add_job(job)
    
    # 触发 worker（异步）
    asyncio.create_task(run_worker(job.id))
    
    return CreateJobResponse(job_id=job_id, status=JobStatus.QUEUED)

@app.get("/v1/jobs/{job_id}", response_model=JobStatusResponse)
async def get_job(job_id: str):
    """查询任务状态"""
    job = store.get_job(job_id)
    if not job:
        raise HTTPException(404, "Job not found")
    
    return JobStatusResponse(
        job_id=job.id,
        status=job.status,
        progress=job.progress,
        stage=job.stage,
        result=job.result,
        quality=job.result.quality if job.result else None,
        error_message=job.error_message
    )

@app.get("/v1/jobs/{job_id}/result")
async def get_job_result(job_id: str):
    """获取任务结果"""
    job = store.get_job(job_id)
    if not job:
        raise HTTPException(404, "Job not found")
    
    if job.status != JobStatus.DONE:
        raise HTTPException(400, f"Job not done yet: {job.status}")
    
    return job.result

@app.get("/v1/jobs/{job_id}/download")
async def download_result(job_id: str):
    """下载 PPTX 文件"""
    job = store.get_job(job_id)
    if not job or not job.result:
        raise HTTPException(404, "Result not found")
    
    # 解析本地路径
    output_url = job.result.output_url
    if output_url.startswith("file://"):
        file_path = output_url[7:]
        return FileResponse(
            file_path,
            media_type="application/vnd.openxmlformats-officedocument.presentationml.presentation",
            filename=f"result_{job_id}.pptx"
        )
    
    raise HTTPException(400, "Invalid output URL")

# ============ Worker (Simulated) ============
async def run_worker(job_id: str):
    """模拟 GPU worker 处理流程"""
    from pptx import Presentation
    from pptx.util import Inches, Pt
    from PIL import Image
    import random
    
    job = store.get_job(job_id)
    if not job:
        return
    
    asset = store.get_asset(job.asset_id)
    
    # 阶段定义（不包括 done，done 在最后设置）
    stages = [
        ("preprocessing", 15),
        ("inferencing", 45),
        ("reconstructing", 75),
        ("scoring", 95),
    ]
    
    try:
        for stage_name, progress in stages:
            # 更新状态
            job.status = JobStatus(stage_name)
            job.stage = stage_name
            job.progress = progress
            store.update_job(job)
            
            logger.info(f"[Worker] Job {job_id}: {stage_name} ({progress}%)")
            
            # 模拟处理时间
            await asyncio.sleep(2)
        
        # 生成示例 PPTX
        prs = Presentation()
        slide_layout = prs.slide_layouts[6]  # blank
        slide = prs.slides.add_slide(slide_layout)
        
        # 添加标题
        title = slide.shapes.add_textbox(Inches(0.5), Inches(0.5), Inches(9), Inches(1))
        tf = title.text_frame
        p = tf.paragraphs[0]
        p.text = "Proposal Diagram - Generated by AI"
        p.font.size = Pt(32)
        p.font.bold = True
        
        # 添加描述
        desc = slide.shapes.add_textbox(Inches(0.5), Inches(2), Inches(9), Inches(4))
        tf = desc.text_frame
        tf.word_wrap = True
        p = tf.paragraphs[0]
        p.text = f"""这是一个由 AI 生成的示例 PPT。

原始文件: {asset.filename if asset else 'unknown'}
任务ID: {job_id}
生成时间: {datetime.now().strftime('%Y-%m-%d %H:%M:%S')}

系统架构图将在这里展示...
- 用户层
- 应用层
- 数据层
- 基础设施层
"""
        p.font.size = Pt(14)
        
        # 保存
        output_path = OUTPUT_DIR / f"{job_id}.pptx"
        prs.save(str(output_path))
        
        # 生成 preview (简单的占位图)
        preview_path = OUTPUT_DIR / f"{job_id}_preview.png"
        img = Image.new('RGB', (800, 600), color=(73, 109, 137))
        img.save(preview_path)
        
        # 质量评分（模拟）
        quality = QualityScore(
            editable_rate=round(random.uniform(0.85, 0.95), 2),
            ocr_accuracy=round(random.uniform(0.92, 0.98), 2),
            layout_deviation=round(random.uniform(0.05, 0.12), 2),
            score=random.randint(82, 95)
        )
        
        # 更新结果
        job.result = JobResult(
            output_url=f"file://{output_path}",
            preview_url=f"file://{preview_path}",
            quality=quality
        )
        job.status = JobStatus.DONE
        job.progress = 100
        job.stage = "done"
        store.update_job(job)
        
        logger.info(f"[Worker] Job {job_id} completed! Quality score: {quality.score}")
        
    except Exception as e:
        logger.error(f"[Worker] Job {job_id} failed: {e}")
        job.status = JobStatus.FAILED_TERMINAL
        job.error_message = str(e)
        store.update_job(job)

# ============ Main ============
if __name__ == "__main__":
    import uvicorn
    uvicorn.run(app, host="0.0.0.0", port=8000, log_level="info")
